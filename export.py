from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from database import supabase
from docx import Document
from pptx import Presentation
import os 
import uuid
import tempfile

def text_to_bullets(text: str, max_lines: int = 6):
    sentences = [s.strip() for s in text.split(".") if s.strip()]

    bullets = []
    for s in sentences:
        if len(s) > 180:
            bullets.append(s[:180] + "…")
        else:
            bullets.append(s)
#final
        if len(bullets) >= max_lines:
            break
    if not bullets:
        return [text[:200] + "…"]

    return bullets

router=APIRouter(prefix="/export",tags=["Export"])

def get_project_info(project_id: str):
    """Fetch project topic and doc_type."""
    project = supabase.table("projects") \
        .select("topic, doc_type") \
        .eq("id", project_id) \
        .execute()
    if project.data:
        return project.data[0].get("topic", "Document"), project.data[0].get("doc_type", "docx")
    return "Document", "docx"

def get_section_titles(project_id: str, doc_type: str):
    """Fetch section/slide titles mapped by ID."""
    if doc_type == "pptx":
        data = supabase.table("slides") \
            .select("id, title") \
            .eq("project_id", project_id) \
            .order("position") \
            .execute()
    else:
        data = supabase.table("sections") \
            .select("id, title") \
            .eq("project_id", project_id) \
            .order("position") \
            .execute()
    
    if data.data:
        return {item["id"]: item["title"] for item in data.data}
    return {}

def get_content(project_id: str):
    ai_data=supabase.table("ai_generations") \
        .select("section_id,content") \
        .eq("project_id",project_id) \
        .execute()
    if not ai_data.data:
        raise HTTPException(status_code=400,detail="no content found for this project")
    content_map={item["section_id"]:{"text":item["content"],"source":"ai"}
    for item in ai_data.data}
    refined=supabase.table("refined_content") \
        .select("section_id,refined_text") \
        .eq("project_id",project_id) \
        .execute()
    if refined.data:
        for item in refined.data:
            content_map[item["section_id"]]={"text":item["refined_text"],"source":"refined"}
    final_content=[{"section_id":k,"text":v["text"],"source":v["source"]}for k,v in content_map.items()]
    return final_content

@router.get("/docx/{project_id}")
def export_docx(project_id: str):
    topic, doc_type = get_project_info(project_id)
    title_map = get_section_titles(project_id, doc_type)
    data = get_content(project_id)

    doc = Document()
    doc.add_heading(topic, level=1)

    for item in data:
        section_title = title_map.get(item["section_id"], "Untitled Section")
        doc.add_heading(section_title, level=2)
        doc.add_paragraph(item["text"])
        doc.add_page_break()

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    filename = tmp.name
    tmp.close()
    doc.save(filename)

    safe_topic = "".join(c if c.isalnum() or c in " _-" else "" for c in topic)[:50]
    return FileResponse(
        filename,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=f"{safe_topic}.docx"
    )

@router.get("/pptx/{project_id}")
def export_pptx(project_id: str):
    topic, doc_type = get_project_info(project_id)
    title_map = get_section_titles(project_id, doc_type)
    data = get_content(project_id)

    prs = Presentation()

    # Title slide with the actual topic
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = "Generated with AI Document Generator"

    for item in data:
        slide_title = title_map.get(item["section_id"], "Untitled Slide")
        bullets = text_to_bullets(item["text"])

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title

        body = slide.placeholders[1].text_frame
        body.text = bullets[0]  
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
            p.level = 0

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    filename = tmp.name
    tmp.close()
    prs.save(filename)

    safe_topic = "".join(c if c.isalnum() or c in " _-" else "" for c in topic)[:50]
    return FileResponse(
        filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{safe_topic}.pptx"
    )
